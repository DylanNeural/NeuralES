"""Génère NeuralES.pptx — importable directement dans Google Slides."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──
BG       = RGBColor(0x05, 0x0A, 0x14)   # fond sombre
ACCENT   = RGBColor(0x00, 0xD4, 0xFF)   # cyan
ACCENT2  = RGBColor(0x00, 0xFF, 0x9D)   # vert
YELLOW   = RGBColor(0xFF, 0xD6, 0x00)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x94, 0xA3, 0xB8)
DARK_BOX = RGBColor(0x0D, 0x1B, 0x2E)
RED      = RGBColor(0xFF, 0x6B, 0x6B)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # blank


def new_slide():
    s = prs.slides.add_slide(blank_layout)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s


def box(slide, x, y, w, h, color=DARK_BOX, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h, size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def accent_bar(slide, y=0.55, color=ACCENT):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(y), Inches(13.33), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def slide_header(slide, title, subtitle=None):
    accent_bar(slide, y=0.55)
    txt(slide, title, 0.5, 0.08, 11, 0.5, size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.5, 0.58, 10, 0.4, size=13, color=GRAY)


def bullet_box(slide, items, x, y, w, h, title=None, icon_color=ACCENT):
    box(slide, x, y, w, h)
    dy = y + 0.15
    if title:
        txt(slide, title, x + 0.2, dy, w - 0.3, 0.35, size=13, bold=True, color=icon_color)
        dy += 0.35
    for item in items:
        txt(slide, f"▸  {item}", x + 0.2, dy, w - 0.3, 0.38, size=12, color=WHITE)
        dy += 0.37


def tag(slide, text, x, y, color=ACCENT):
    b = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.2), Inches(0.38))
    b.fill.solid()
    b.fill.fore_color.rgb = color
    b.line.fill.background()
    tf = b.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = BG


# ══════════════════════════════════════════════
# SLIDE 1 — Titre
# ══════════════════════════════════════════════
s = new_slide()
accent_bar(s, y=3.5, color=ACCENT)
accent_bar(s, y=3.56, color=ACCENT2)

txt(s, "NeuralES", 1.5, 1.0, 10, 1.5, size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Système d'analyse EEG de la fatigue cognitive en temps réel",
    1.0, 2.5, 11.33, 0.8, size=20, color=GRAY, align=PP_ALIGN.CENTER)

txt(s, "Dylan Andrade Pereira  ·  Ynov Campus Toulouse  ·  Juin 2026",
    1.0, 6.5, 11.33, 0.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)

# déco
for i, col in enumerate([ACCENT, ACCENT2, YELLOW, RED]):
    b = s.shapes.add_shape(1, Inches(0.3 + i * 0.25), Inches(4.0), Inches(0.08), Inches(2.5))
    b.fill.solid()
    b.fill.fore_color.rgb = col
    b.line.fill.background()

# ══════════════════════════════════════════════
# SLIDE 2 — Sommaire
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Sommaire")

items = [
    ("01", "Contexte & Problématique",  ACCENT),
    ("02", "Objectifs du projet",        ACCENT2),
    ("03", "Architecture globale",       YELLOW),
    ("04", "Réalisations techniques",    RGBColor(0xC0,0x84,0xFC)),
    ("05", "Démonstration",              RGBColor(0xFB,0x92,0x3C)),
    ("06", "Difficultés rencontrées",    RED),
    ("07", "Perspectives & Conclusion",  ACCENT),
]

for i, (num, label, col) in enumerate(items):
    row = i % 4
    col_n = i // 4
    bx = 0.5 + col_n * 6.8
    by = 1.2 + row * 1.4
    bw = 6.3
    bh = 1.2
    b = s.shapes.add_shape(1, Inches(bx), Inches(by), Inches(bw), Inches(bh))
    b.fill.solid()
    b.fill.fore_color.rgb = DARK_BOX
    b.line.fill.background()
    accent = s.shapes.add_shape(1, Inches(bx), Inches(by), Inches(0.06), Inches(bh))
    accent.fill.solid()
    accent.fill.fore_color.rgb = col
    accent.line.fill.background()
    txt(s, num, bx + 0.15, by + 0.15, 0.6, 0.5, size=20, bold=True, color=col)
    txt(s, label, bx + 0.75, by + 0.28, bw - 0.9, 0.5, size=15, bold=True, color=WHITE)

# ══════════════════════════════════════════════
# SLIDE 3 — Contexte : fatigue cognitive
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "La fatigue cognitive : un enjeu invisible", "Contexte")

box(s, 0.4, 1.2, 5.8, 5.6)
txt(s, "Définition", 0.7, 1.35, 5.2, 0.4, size=14, bold=True, color=ACCENT)
txt(s, "Dégradation des performances mentales due à une activité intellectuelle soutenue : baisse de concentration, augmentation du temps de réaction, mauvaise prise de décision.",
    0.7, 1.75, 5.2, 1.2, size=12, color=WHITE)

txt(s, "Domaines à risque", 0.7, 3.0, 5.2, 0.4, size=14, bold=True, color=ACCENT)
for item in ["Chirurgie & médecine", "Aviation & transport", "Industrie & nucléaire", "Sport de haut niveau"]:
    txt(s, f"▸  {item}", 0.7, 3.4 + ["Chirurgie & médecine", "Aviation & transport", "Industrie & nucléaire", "Sport de haut niveau"].index(item) * 0.42, 5.2, 0.4, size=12, color=WHITE)

txt(s, "Problème actuel", 0.7, 5.55, 5.2, 0.4, size=14, bold=True, color=RED)
txt(s, "Mesure subjective uniquement (questionnaires) — pas de détection objective en temps réel.",
    0.7, 5.95, 5.2, 0.7, size=12, color=WHITE)

box(s, 6.6, 1.2, 6.3, 2.5)
txt(s, "Impact réel", 6.9, 1.35, 5.7, 0.4, size=14, bold=True, color=YELLOW)
txt(s, "Un chirurgien après 24h de garde commet 20% d'erreurs supplémentaires (étude Harvard, 2004)",
    6.9, 1.8, 5.7, 1.6, size=13, color=WHITE)

box(s, 6.6, 4.0, 6.3, 2.8)
txt(s, "Solution envisagée", 6.9, 4.15, 5.7, 0.4, size=14, bold=True, color=ACCENT2)
txt(s, "L'EEG permet de détecter des marqueurs biomédicaux objectifs de la fatigue cognitive en temps réel.",
    6.9, 4.6, 5.7, 1.8, size=13, color=WHITE)

# ══════════════════════════════════════════════
# SLIDE 4 — L'EEG comme solution
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "L'EEG comme outil de mesure", "Contexte")

bands = [
    ("Delta",  "0.5–4 Hz",  "Sommeil profond",       GRAY,    0.15),
    ("Thêta",  "4–8 Hz",    "Somnolence — AUGMENTE avec fatigue", RED,    0.55),
    ("Alpha",  "8–13 Hz",   "Relaxation — DIMINUE avec vigilance", YELLOW, 0.45),
    ("Bêta",   "13–30 Hz",  "Concentration active",  ACCENT2, 0.35),
    ("Gamma",  "> 30 Hz",   "Traitement cognitif",   ACCENT,  0.25),
]
for i, (name, freq, desc, col, _) in enumerate(bands):
    bx, by = 0.4, 1.2 + i * 1.1
    b = s.shapes.add_shape(1, Inches(bx), Inches(by), Inches(8.5), Inches(0.9))
    b.fill.solid(); b.fill.fore_color.rgb = DARK_BOX; b.line.fill.background()
    accent_s = s.shapes.add_shape(1, Inches(bx), Inches(by), Inches(0.06), Inches(0.9))
    accent_s.fill.solid(); accent_s.fill.fore_color.rgb = col; accent_s.line.fill.background()
    txt(s, name, bx+0.15, by+0.1, 1.2, 0.4, size=14, bold=True, color=col)
    txt(s, freq, bx+1.4, by+0.1, 1.5, 0.4, size=12, color=GRAY)
    txt(s, desc, bx+3.0, by+0.1, 5.5, 0.4, size=12, color=WHITE)

box(s, 9.3, 1.2, 3.7, 5.5)
txt(s, "Marqueur clé", 9.55, 1.4, 3.2, 0.4, size=13, bold=True, color=ACCENT)
txt(s, "Indice de fatigue :", 9.55, 1.85, 3.2, 0.4, size=12, color=WHITE)

b2 = s.shapes.add_shape(1, Inches(9.55), Inches(2.4), Inches(3.2), Inches(1.1))
b2.fill.solid(); b2.fill.fore_color.rgb = RGBColor(0x0D,0x2B,0x1A); b2.line.fill.background()
txt(s, "Thêta / Alpha", 9.65, 2.55, 3.0, 0.7, size=22, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

txt(s, "Un ratio élevé = fatigue détectée. Mesurable en continu, en temps réel.",
    9.55, 3.7, 3.2, 1.2, size=11, color=WHITE)

txt(s, "Matériel : OpenBCI Ultracortex", 9.55, 5.1, 3.2, 0.4, size=12, bold=True, color=YELLOW)
txt(s, "8–16 canaux · 250 Hz · Open source · ~300€",
    9.55, 5.5, 3.2, 0.8, size=11, color=WHITE)

# ══════════════════════════════════════════════
# SLIDE 5 — Problématique
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Problématique", "")

b = s.shapes.add_shape(1, Inches(0.9), Inches(1.3), Inches(11.5), Inches(2.2))
b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x03,0x14,0x26); b.line.fill.background()
a = s.shapes.add_shape(1, Inches(0.9), Inches(1.3), Inches(0.08), Inches(2.2))
a.fill.solid(); a.fill.fore_color.rgb = ACCENT; a.line.fill.background()
txt(s, "Comment concevoir une application desktop capable d'acquérir, visualiser et analyser des signaux EEG en temps réel, adaptée à un usage clinique ?",
    1.2, 1.5, 11.0, 1.8, size=18, color=WHITE)

txt(s, "Contraintes", 0.9, 3.9, 4.0, 0.4, size=15, bold=True, color=ACCENT)
for c in ["Traitement signal bas-latence (< 200ms)", "Interface lisible par un non-technicien", "Architecture extensible (multi-patients, historique)", "Windows 11 — déploiement simple"]:
    txt(s, f"▸  {c}", 0.9, 4.3 + ["Traitement signal bas-latence (< 200ms)", "Interface lisible par un non-technicien", "Architecture extensible (multi-patients, historique)", "Windows 11 — déploiement simple"].index(c)*0.48, 5.5, 0.45, size=13, color=WHITE)

txt(s, "Périmètre fil rouge", 7.0, 3.9, 5.0, 0.4, size=15, bold=True, color=ACCENT2)
txt(s, "Phase 1 (ce semestre) :", 7.0, 4.4, 5.0, 0.35, size=13, bold=True, color=WHITE)
txt(s, "Stack complète backend + desktop fonctionnelle, données EEG via fichier EDF (PhysioNet Sleep)",
    7.0, 4.75, 5.8, 0.9, size=12, color=WHITE)
txt(s, "Phase 2 :", 7.0, 5.7, 5.0, 0.35, size=13, bold=True, color=GRAY)
txt(s, "Connexion matériel OpenBCI réel + algorithme fatigue",
    7.0, 6.05, 5.8, 0.6, size=12, color=GRAY)

# ══════════════════════════════════════════════
# SLIDE 6 — Objectifs
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Objectifs du projet", "")

goals = [
    ("Application desktop PySide6 complète",    True),
    ("Gestion des patients (CRUD)",              True),
    ("Acquisition EEG temps réel (WebSocket)",   True),
    ("Visualisation signal waveform",            True),
    ("Graphique 3D waterfall (Three.js)",        True),
    ("Casque EEG 3D interactif",                 True),
    ("Backend FastAPI + BDD PostgreSQL (VPS)",   True),
    ("Connexion matériel OpenBCI",               False),
]

for i, (goal, done) in enumerate(goals):
    col_n = i // 4
    row   = i % 4
    bx = 0.4 + col_n * 6.55
    by = 1.2 + row * 1.4
    b = s.shapes.add_shape(1, Inches(bx), Inches(by), Inches(6.1), Inches(1.15))
    b.fill.solid(); b.fill.fore_color.rgb = DARK_BOX; b.line.fill.background()
    col = ACCENT2 if done else YELLOW
    ic = "✅" if done else "🔄"
    txt(s, ic, bx+0.15, by+0.25, 0.6, 0.5, size=20, color=col)
    txt(s, goal, bx+0.7, by+0.28, 5.2, 0.55, size=13, bold=done, color=WHITE)
    status = "Réalisé" if done else "En cours"
    tag_b = s.shapes.add_shape(1, Inches(bx+0.7), Inches(by+0.72), Inches(1.4), Inches(0.28))
    tag_b.fill.solid(); tag_b.fill.fore_color.rgb = col if not done else RGBColor(0x05,0x2A,0x1A); tag_b.line.fill.background()
    tf = tag_b.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = status
    r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = col if done else BG

# ══════════════════════════════════════════════
# SLIDE 7 — Architecture
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Architecture Globale", "")

# ── Couche 1 : Desktop ──
desktop_y = 0.9
b = s.shapes.add_shape(1, Inches(0.5), Inches(desktop_y), Inches(12.3), Inches(2.2))
b.fill.solid(); b.fill.fore_color.rgb = DARK_BOX; b.line.fill.background()
top = s.shapes.add_shape(1, Inches(0.5), Inches(desktop_y), Inches(12.3), Inches(0.06))
top.fill.solid(); top.fill.fore_color.rgb = ACCENT; top.line.fill.background()
txt(s, "APPLICATION DESKTOP  —  PySide6 / Qt 6", 0.7, desktop_y+0.08, 12.0, 0.35, size=12, bold=True, color=ACCENT)

for i, (name, col) in enumerate([("Patients (CRUD)", ACCENT), ("Acquisition + Signaux", ACCENT2), ("Brain View (Three.js)", YELLOW)]):
    bx = 0.7 + i * 4.1
    b2 = s.shapes.add_shape(1, Inches(bx), Inches(desktop_y+0.55), Inches(3.8), Inches(1.38))
    b2.fill.solid(); b2.fill.fore_color.rgb = RGBColor(0x0D,0x20,0x35); b2.line.fill.background()
    txt(s, name, bx+0.1, desktop_y+0.92, 3.6, 0.5, size=13, color=col, align=PP_ALIGN.CENTER)

# ── Flèche 1 : Desktop → Backend ──
arrow1_y = 3.25
line1 = s.shapes.add_shape(1, Inches(6.5), Inches(arrow1_y), Inches(0.06), Inches(0.55))
line1.fill.solid(); line1.fill.fore_color.rgb = GRAY; line1.line.fill.background()
tri1 = s.shapes.add_shape(5, Inches(6.28), Inches(arrow1_y+0.55), Inches(0.5), Inches(0.28))
tri1.fill.solid(); tri1.fill.fore_color.rgb = GRAY; tri1.line.fill.background()
txt(s, "HTTP REST  +  WebSocket", 3.5, arrow1_y+0.1, 5.5, 0.4, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ── Couche 2 : Backend ──
backend_y = 4.1
b = s.shapes.add_shape(1, Inches(0.5), Inches(backend_y), Inches(12.3), Inches(1.8))
b.fill.solid(); b.fill.fore_color.rgb = DARK_BOX; b.line.fill.background()
top = s.shapes.add_shape(1, Inches(0.5), Inches(backend_y), Inches(12.3), Inches(0.06))
top.fill.solid(); top.fill.fore_color.rgb = ACCENT2; top.line.fill.background()
txt(s, "BACKEND  —  FastAPI + Uvicorn", 0.7, backend_y+0.08, 12.0, 0.35, size=12, bold=True, color=ACCENT2)

for i, (name, col) in enumerate([("/patients  (CRUD)", ACCENT), ("/acquisition  start/stop", ACCENT2), ("/eeg/stream  WebSocket", YELLOW)]):
    bx = 0.7 + i * 4.1
    b2 = s.shapes.add_shape(1, Inches(bx), Inches(backend_y+0.52), Inches(3.8), Inches(1.0))
    b2.fill.solid(); b2.fill.fore_color.rgb = RGBColor(0x0D,0x20,0x35); b2.line.fill.background()
    txt(s, name, bx+0.1, backend_y+0.75, 3.6, 0.45, size=12, color=col, align=PP_ALIGN.CENTER)

# ── Flèche 2 : Backend → PostgreSQL ──
arrow2_y = 6.05
line2 = s.shapes.add_shape(1, Inches(6.5), Inches(arrow2_y), Inches(0.06), Inches(0.45))
line2.fill.solid(); line2.fill.fore_color.rgb = GRAY; line2.line.fill.background()
tri2 = s.shapes.add_shape(5, Inches(6.28), Inches(arrow2_y+0.45), Inches(0.5), Inches(0.25))
tri2.fill.solid(); tri2.fill.fore_color.rgb = GRAY; tri2.line.fill.background()
txt(s, "SSH Tunnel  localhost:5433 -> VPS:5432", 3.5, arrow2_y+0.05, 5.5, 0.4, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ── Couche 3 : PostgreSQL ──
pg_y = 6.75
b = s.shapes.add_shape(1, Inches(0.5), Inches(pg_y), Inches(12.3), Inches(0.7))
b.fill.solid(); b.fill.fore_color.rgb = DARK_BOX; b.line.fill.background()
top = s.shapes.add_shape(1, Inches(0.5), Inches(pg_y), Inches(12.3), Inches(0.06))
top.fill.solid(); top.fill.fore_color.rgb = YELLOW; top.line.fill.background()
txt(s, "PostgreSQL 15  —  VPS OVH Debian  (51.178.30.35)", 0.7, pg_y+0.15, 12.0, 0.4, size=12, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 8 — Stack Technique
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Stack Technique", "")

stacks = [
    ("Backend", [
        "Python 3.11",
        "FastAPI + Uvicorn (async)",
        "SQLAlchemy 2.0 + Alembic",
        "PostgreSQL 15",
        "MNE (lecture EDF)",
        "WebSocket natif asyncio",
    ], ACCENT2),
    ("Desktop", [
        "PySide6 / Qt 6.6",
        "pyqtgraph (waveform)",
        "QWebEngineView (Three.js)",
        "QWebChannel (Python↔JS)",
        "NumPy (buffer circulaire)",
        "Three.js r168 (3D)",
    ], ACCENT),
    ("Infrastructure", [
        "VPS OVH — Debian 12",
        "PostgreSQL 15 distant",
        "Tunnel SSH local :5433",
        "Serveur HTTP embarqué",
        "Données : Sleep-EDF (PhysioNet)",
        "OS : Windows 11 Pro",
    ], YELLOW),
]

for i, (title, items, col) in enumerate(stacks):
    bx = 0.4 + i * 4.3
    b = s.shapes.add_shape(1, Inches(bx), Inches(1.2), Inches(4.0), Inches(5.8))
    b.fill.solid(); b.fill.fore_color.rgb = DARK_BOX; b.line.fill.background()
    top = s.shapes.add_shape(1, Inches(bx), Inches(1.2), Inches(4.0), Inches(0.06))
    top.fill.solid(); top.fill.fore_color.rgb = col; top.line.fill.background()
    txt(s, title, bx+0.2, 1.28, 3.6, 0.45, size=15, bold=True, color=col)
    for j, item in enumerate(items):
        txt(s, f"▸  {item}", bx+0.2, 1.85+j*0.72, 3.6, 0.5, size=12, color=WHITE)

# ══════════════════════════════════════════════
# SLIDE 9 — Gestion des patients
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Gestion des Patients", "Réalisations")

box(s, 0.4, 1.2, 6.0, 5.8)
txt(s, "Liste patients", 0.7, 1.35, 5.5, 0.4, size=14, bold=True, color=ACCENT)
for item in [
    "Avatars colorés (initiales + couleur par hash)",
    "Badges service colorés",
    "4 KPI cards (stats globales)",
    "Filtre temps réel : nom, service, médecin",
    "Bouton 'Voir' par ligne",
]:
    txt(s, f"▸  {item}", 0.7, 1.8+["Avatars colorés (initiales + couleur par hash)","Badges service colorés","4 KPI cards (stats globales)","Filtre temps réel : nom, service, médecin","Bouton 'Voir' par ligne"].index(item)*0.55, 5.5, 0.5, size=12, color=WHITE)

txt(s, "Création de dossier", 0.7, 4.7, 5.5, 0.4, size=14, bold=True, color=ACCENT2)
for item in [
    "Nom, Prénom, Date de naissance, Sexe, NSS",
    "Service, Médecin référent, Notes, Remarques",
    "Scroll complet jusqu'en bas",
]:
    txt(s, f"▸  {item}", 0.7, 5.2+["Nom, Prénom, Date de naissance, Sexe, NSS","Service, Médecin référent, Notes, Remarques","Scroll complet jusqu'en bas"].index(item)*0.52, 5.5, 0.45, size=12, color=WHITE)

# Mockup patients list
box(s, 6.8, 1.2, 6.1, 5.8, color=RGBColor(0x08,0x12,0x1F))
txt(s, "[ Screenshot liste patients ]", 7.0, 3.7, 5.7, 0.6, size=14, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "Insérer capture d'écran ici", 7.0, 4.3, 5.7, 0.4, size=11, color=RGBColor(0x3B,0x4E,0x63), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 10 — Acquisition EEG
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Acquisition EEG — Architecture Temps Réel", "Réalisations")

steps = [
    ("1", "POST /acquisition/start", "Crée une session en BDD\nRetourne session_id", ACCENT),
    ("2", "WS /eeg/stream?session_id=…", "Connexion WebSocket\nStream JSON chunks 100ms", ACCENT2),
    ("3", "_on_ws_msg()", "Parse chunk Python\nMet à jour les graphiques", YELLOW),
]

for i, (num, title, desc, col) in enumerate(steps):
    bx = 0.5 + i * 4.2
    b = s.shapes.add_shape(1, Inches(bx), Inches(1.3), Inches(3.8), Inches(2.6))
    b.fill.solid(); b.fill.fore_color.rgb = DARK_BOX; b.line.fill.background()
    circle = s.shapes.add_shape(9, Inches(bx+1.55), Inches(1.4), Inches(0.7), Inches(0.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = col; circle.line.fill.background()
    txt(s, num, bx+1.55, 1.4, 0.7, 0.7, size=18, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, title, bx+0.1, 2.25, 3.6, 0.5, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, desc, bx+0.1, 2.8, 3.6, 0.9, size=11, color=WHITE, align=PP_ALIGN.CENTER)

    if i < 2:
        arr = s.shapes.add_shape(1, Inches(bx+3.8), Inches(2.4), Inches(0.4), Inches(0.2))
        arr.fill.solid(); arr.fill.fore_color.rgb = GRAY; arr.line.fill.background()

box(s, 0.4, 4.2, 12.5, 2.9)
txt(s, "Format d'un chunk WebSocket", 0.7, 4.35, 12.0, 0.4, size=13, bold=True, color=ACCENT)
code = '{\n  "channels": ["EEG Fpz-Cz", "EEG Pz-Oz", ...],\n  "samples":  [[-0.000045, 0.000023, ...], [...]],\n  "t0":       142.3,          // timestamp départ chunk (secondes)\n  "sfreq":    100.0           // fréquence d\'échantillonnage (Hz)\n}'
txt(s, code, 0.7, 4.8, 12.0, 2.1, size=11, color=ACCENT2)

# ══════════════════════════════════════════════
# SLIDE 11 — Visualisation waveform
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Visualisation Signal — Waveform (pyqtgraph)", "Réalisations")

box(s, 0.4, 1.2, 5.5, 5.8)
txt(s, "Fonctionnement", 0.7, 1.35, 5.0, 0.4, size=14, bold=True, color=ACCENT)
for item in [
    "PlotWidget pyqtgraph, courbes multi-canaux",
    "Buffer circulaire NumPy (5000 pts ≈ 20s)",
    "Mise à jour à chaque chunk WebSocket",
    "Défilement temps réel — axe X en secondes",
    "Couleur différente par canal EEG",
]:
    txt(s, f"▸  {item}", 0.7, 1.85+["PlotWidget pyqtgraph, courbes multi-canaux","Buffer circulaire NumPy (5000 pts ≈ 20s)","Mise à jour à chaque chunk WebSocket","Défilement temps réel — axe X en secondes","Couleur différente par canal EEG"].index(item)*0.6, 5.0, 0.5, size=12, color=WHITE)

txt(s, "Problème résolu", 0.7, 5.0, 5.0, 0.4, size=13, bold=True, color=RED)
txt(s, "Conflit D3D11 (Qt) / OpenGL (pyqtgraph) sur Windows 11", 0.7, 5.45, 5.0, 0.5, size=12, color=WHITE)
code2 = "pg.setConfigOptions(useOpenGL=False)"
txt(s, code2, 0.7, 6.0, 5.0, 0.5, size=11, color=ACCENT2)

box(s, 6.3, 1.2, 6.6, 5.8, color=RGBColor(0x08,0x12,0x1F))
txt(s, "[ Screenshot waveform ]", 6.5, 3.7, 6.2, 0.6, size=14, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "Insérer capture d'écran ici", 6.5, 4.3, 6.2, 0.4, size=11, color=RGBColor(0x3B,0x4E,0x63), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 12 — Graphique 3D
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Visualisation 3D — Waterfall Chart (Three.js)", "Réalisations")

box(s, 0.4, 1.2, 5.5, 5.8)
txt(s, "Concept", 0.7, 1.35, 5.0, 0.4, size=14, bold=True, color=ACCENT)
txt(s, "Chaque canal EEG = une courbe 3D. Les canaux sont empilés en profondeur (axe Z). Le temps défile sur l'axe X (fenêtre 10s).",
    0.7, 1.8, 5.0, 1.3, size=12, color=WHITE)

txt(s, "Implémentation", 0.7, 3.2, 5.0, 0.4, size=14, bold=True, color=ACCENT2)
for item in [
    "THREE.BufferGeometry + Float32Array",
    "Buffer circulaire JS (800 points / canal)",
    "Auto-scaling au 1er chunk (Volts → scène)",
    "Fog exponentiel, grille de fond",
    "OrbitControls (rotation / zoom / pan)",
]:
    txt(s, f"▸  {item}", 0.7, 3.7+["THREE.BufferGeometry + Float32Array","Buffer circulaire JS (800 points / canal)","Auto-scaling au 1er chunk (Volts → scène)","Fog exponentiel, grille de fond","OrbitControls (rotation / zoom / pan)"].index(item)*0.5, 5.0, 0.45, size=12, color=WHITE)

txt(s, "Auto-scaling", 0.7, 6.2, 5.0, 0.35, size=12, bold=True, color=YELLOW)
txt(s, "ampScale = 55 / maxAbsValue", 0.7, 6.55, 5.0, 0.35, size=11, color=ACCENT2)

box(s, 6.3, 1.2, 6.6, 5.8, color=RGBColor(0x08,0x12,0x1F))
txt(s, "[ Screenshot graphique 3D ]", 6.5, 3.7, 6.2, 0.6, size=14, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "Insérer capture d'écran ici", 6.5, 4.3, 6.2, 0.4, size=11, color=RGBColor(0x3B,0x4E,0x63), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 13 — Casque EEG 3D
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Casque EEG 3D Interactif", "Réalisations")

box(s, 0.4, 1.2, 5.5, 5.8)
txt(s, "Fonctionnalités", 0.7, 1.35, 5.0, 0.4, size=14, bold=True, color=ACCENT)
for item in [
    "Modèle GLB avec électrodes nommées",
    "Clic sur électrode → sélection",
    "Coloration par score de fatigue",
    "Rotation / zoom via OrbitControls",
    "Chargement via GLTFLoader",
]:
    txt(s, f"▸  {item}", 0.7, 1.85+["Modèle GLB avec électrodes nommées","Clic sur électrode → sélection","Coloration par score de fatigue","Rotation / zoom via OrbitControls","Chargement via GLTFLoader"].index(item)*0.6, 5.0, 0.5, size=12, color=WHITE)

txt(s, "Communication Python ↔ JS", 0.7, 5.05, 5.0, 0.4, size=13, bold=True, color=ACCENT2)
code3 = "# Python reçoit :\ndef on_electrode_click(self, json_str):\n    data = json.loads(json_str)\n\n// JS envoie :\n_bridge.on_electrode_click(\n    JSON.stringify({electrode, active})\n);"
txt(s, code3, 0.7, 5.5, 5.0, 1.6, size=10, color=ACCENT2)

box(s, 6.3, 1.2, 6.6, 5.8, color=RGBColor(0x08,0x12,0x1F))
txt(s, "[ Screenshot casque 3D ]", 6.5, 3.7, 6.2, 0.6, size=14, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "Insérer capture d'écran ici", 6.5, 4.3, 6.2, 0.4, size=11, color=RGBColor(0x3B,0x4E,0x63), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 14 — QWebChannel
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Bridge Python ↔ JavaScript (QWebChannel)", "Réalisations")

box(s, 0.4, 1.2, 5.8, 5.8)
txt(s, "Côté Python (PySide6)", 0.7, 1.35, 5.3, 0.4, size=13, bold=True, color=ACCENT)
code_py = """from PySide6.QtCore import QObject, Slot
from PySide6.QtWebChannel import QWebChannel

class Bridge(QObject):
    @Slot(str)
    def on_electrode_click(self, json_str):
        data = json.loads(json_str)
        # Traitement de la sélection

channel = QWebChannel()
channel.registerObject('bridge', bridge)
page.setWebChannel(channel)"""
txt(s, code_py, 0.7, 1.85, 5.3, 3.5, size=10, color=ACCENT2)

txt(s, "Côté JavaScript (Three.js)", 0.7, 5.5, 5.3, 0.4, size=13, bold=True, color=YELLOW)
code_js = "new QWebChannel(qt.webChannelTransport,\n  ch => { window._bridge = ch.objects.bridge; });\n\n// Envoyer vers Python :\n_bridge.on_electrode_click(JSON.stringify({electrode, active}));"
txt(s, code_js, 0.7, 5.95, 5.3, 1.3, size=10, color=YELLOW)

box(s, 6.6, 1.2, 6.3, 5.8)
txt(s, "Pourquoi QWebChannel ?", 6.9, 1.35, 5.8, 0.4, size=13, bold=True, color=ACCENT)
for item in [
    "Communication directe Python ↔ JS",
    "Pas de serveur HTTP intermédiaire",
    "Synchrone, sans WebSocket externe",
    "Signal/Slot natif Qt",
]:
    txt(s, f"▸  {item}", 6.9, 1.9+["Communication directe Python ↔ JS","Pas de serveur HTTP intermédiaire","Synchrone, sans WebSocket externe","Signal/Slot natif Qt"].index(item)*0.6, 5.8, 0.5, size=12, color=WHITE)

box(s, 6.6, 4.3, 6.3, 2.7)
txt(s, "Serveur HTTP embarqué", 6.9, 4.45, 5.8, 0.4, size=13, bold=True, color=ACCENT2)
txt(s, "Three.js ne peut pas être chargé en file://. Un TCPServer Python démarre sur un port aléatoire au lancement et sert les fichiers vendor + HTML.",
    6.9, 4.9, 5.8, 1.8, size=12, color=WHITE)

# ══════════════════════════════════════════════
# SLIDE 15 — Difficultés
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Difficultés Techniques & Solutions", "")

problems = [
    ("Conflit D3D11 / OpenGL",         "Qt WebEngine + pyqtgraph crash au démarrage",      "pg.setConfigOptions(useOpenGL=False)",        RED,    ACCENT2),
    ("Three.js CDN inaccessible",      "CDN bloqués dans Qt WebEngine, file:// trop limité","Vendor local + serveur HTTP Python embarqué", YELLOW, ACCENT2),
    ("Échelle signaux EEG",            "Valeurs en Volts (±0.0001V) — courbes invisibles",  "Auto-scaling: ampScale = 55 / maxAbsValue",   YELLOW, ACCENT2),
    ("QTextEdit vole le scroll",       "Molette sur textarea → form ne scroll plus",        "wheelEvent = lambda e: e.ignore()",           RED,    ACCENT2),
    ("Session WebSocket requise",      "WS refusé sans session_id valide",                  "Flux HTTP POST → session_id → ouverture WS",  YELLOW, ACCENT2),
    ("VPS PostgreSQL down",            "Backend ne peut pas se connecter à la BDD",         "Redémarrage depuis panel admin OVH",          RED,    ACCENT),
]

for i, (title, prob, sol, pc, sc) in enumerate(problems):
    row = i % 3; col_n = i // 3
    bx = 0.4 + col_n * 6.5
    by = 1.2 + row * 2.0
    b = s.shapes.add_shape(1, Inches(bx), Inches(by), Inches(6.0), Inches(1.75))
    b.fill.solid(); b.fill.fore_color.rgb = DARK_BOX; b.line.fill.background()
    txt(s, title, bx+0.15, by+0.08, 5.7, 0.4, size=12, bold=True, color=WHITE)
    txt(s, f"✗  {prob}", bx+0.15, by+0.5, 5.7, 0.4, size=10, color=pc)
    txt(s, f"✓  {sol}", bx+0.15, by+0.95, 5.7, 0.55, size=10, color=ACCENT2)

# ══════════════════════════════════════════════
# SLIDE 16 — Demo
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Démonstration", "")

b = s.shapes.add_shape(1, Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.8))
b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x08,0x12,0x1F); b.line.fill.background()

txt(s, "LIVE", 5.5, 2.0, 2.3, 0.9, size=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, "DEMO", 5.5, 2.9, 2.3, 0.9, size=48, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

steps_demo = [
    "Lancer le backend (FastAPI + tunnel SSH)",
    "Ouvrir l'application desktop",
    "Créer un patient",
    "Démarrer une session d'acquisition",
    "Visualisation waveform temps réel",
    "Basculer vers graphique 3D waterfall",
    "Interagir avec le casque EEG 3D",
]
for i, step in enumerate(steps_demo):
    bx = 8.2 if i < 4 else 8.2
    by_offset = 1.35 + i * 0.62
    circle = s.shapes.add_shape(9, Inches(8.1), Inches(by_offset), Inches(0.36), Inches(0.36))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT if i % 2 == 0 else ACCENT2; circle.line.fill.background()
    txt(s, str(i+1), 8.1, by_offset, 0.36, 0.36, size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, step, 8.6, by_offset+0.01, 4.1, 0.38, size=11, color=WHITE)

# ══════════════════════════════════════════════
# SLIDE 17 — Résultats
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Résultats", "")

metrics = [
    ("< 100ms",     "Latence WebSocket\n→ rendu graphique",  ACCENT),
    ("7+",          "Canaux EEG\nsimultanés",                ACCENT2),
    ("4",           "Modules\nfonctionnels",                  YELLOW),
    ("100%",        "Stack opérationnelle\nde bout en bout",  RGBColor(0xC0,0x84,0xFC)),
]
for i, (val, label, col) in enumerate(metrics):
    bx = 0.5 + i * 3.1
    b = s.shapes.add_shape(1, Inches(bx), Inches(1.2), Inches(2.8), Inches(2.8))
    b.fill.solid(); b.fill.fore_color.rgb = DARK_BOX; b.line.fill.background()
    top = s.shapes.add_shape(1, Inches(bx), Inches(1.2), Inches(2.8), Inches(0.06))
    top.fill.solid(); top.fill.fore_color.rgb = col; top.line.fill.background()
    txt(s, val, bx, 1.5, 2.8, 1.0, size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, label, bx+0.1, 2.6, 2.6, 1.0, size=12, color=WHITE, align=PP_ALIGN.CENTER)

box(s, 0.4, 4.3, 12.5, 2.7)
txt(s, "Stack complète validée sur Windows 11", 0.7, 4.45, 12.0, 0.4, size=14, bold=True, color=WHITE)
checks = [
    ("Application lance sans erreur", ACCENT2),
    ("Patients créés et persistés en BDD (VPS OVH)", ACCENT2),
    ("Signal EEG reçu et affiché en temps réel", ACCENT2),
    ("Graphique 3D waterfall animé", ACCENT2),
    ("Casque 3D interactif + QWebChannel opérationnel", ACCENT2),
]
for i, (item, col) in enumerate(checks):
    txt(s, f"✅  {item}", 0.7 + (i//3)*6.5, 4.9 + (i%3)*0.55, 6.0, 0.5, size=12, color=WHITE)

# ══════════════════════════════════════════════
# SLIDE 18 — Perspectives
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Perspectives d'Évolution", "")

phases = [
    ("Court terme", [
        "Connexion OpenBCI réel (BrainFlow SDK)",
        "Algorithme fatigue : ratio Thêta / Alpha",
        "FFT temps réel sur fenêtre glissante 4s",
        "Export sessions CSV / EDF",
    ], ACCENT, 0.4),
    ("Moyen terme", [
        "Application mobile (React Native)",
        "Dashboard web Vue 3 pour médecins",
        "Authentification JWT + gestion des rôles",
        "Alertes temps réel (seuils configurables)",
    ], ACCENT2, 4.6),
    ("Long terme", [
        "Modèle ML (LSTM / Transformer sur EEG)",
        "Certification dispositif médical (CE)",
        "Multi-patients simultanés",
        "SaaS cloud pour hôpitaux",
    ], YELLOW, 8.8),
]

for title, items, col, bx in phases:
    b = s.shapes.add_shape(1, Inches(bx), Inches(1.2), Inches(4.1), Inches(5.8))
    b.fill.solid(); b.fill.fore_color.rgb = DARK_BOX; b.line.fill.background()
    top = s.shapes.add_shape(1, Inches(bx), Inches(1.2), Inches(4.1), Inches(0.06))
    top.fill.solid(); top.fill.fore_color.rgb = col; top.line.fill.background()
    txt(s, title, bx+0.15, 1.3, 3.8, 0.45, size=14, bold=True, color=col)
    for j, item in enumerate(items):
        txt(s, f"▸  {item}", bx+0.15, 1.9+j*0.72, 3.8, 0.6, size=12, color=WHITE)

# ══════════════════════════════════════════════
# SLIDE 19 — Conclusion
# ══════════════════════════════════════════════
s = new_slide()
slide_header(s, "Conclusion", "")

box(s, 0.4, 1.2, 12.5, 2.4)
txt(s, "NeuralES démontre la faisabilité technique d'un outil EEG desktop moderne avec une stack Python/Qt performante pour du traitement signal temps réel.",
    0.7, 1.45, 12.0, 1.0, size=15, color=WHITE)
txt(s, "L'architecture modulaire (backend FastAPI indépendant, desktop PySide6, rendu 3D Three.js) offre une base solide pour l'intégration d'un algorithme de détection de fatigue réel.",
    0.7, 2.4, 12.0, 1.0, size=15, color=WHITE)

points = [
    ("Stack bout-en-bout", "FastAPI + PostgreSQL + PySide6 + Three.js", ACCENT),
    ("Défis techniques résolus", "D3D11, WebGL dans Qt, auto-scaling EEG", ACCENT2),
    ("Extensible", "OpenBCI, ML, mobile, web prêts à intégrer", YELLOW),
]
for i, (title, desc, col) in enumerate(points):
    bx = 0.4 + i * 4.3
    b = s.shapes.add_shape(1, Inches(bx), Inches(3.9), Inches(4.0), Inches(1.6))
    b.fill.solid(); b.fill.fore_color.rgb = DARK_BOX; b.line.fill.background()
    left = s.shapes.add_shape(1, Inches(bx), Inches(3.9), Inches(0.06), Inches(1.6))
    left.fill.solid(); left.fill.fore_color.rgb = col; left.line.fill.background()
    txt(s, title, bx+0.15, 3.95, 3.7, 0.4, size=13, bold=True, color=col)
    txt(s, desc, bx+0.15, 4.4, 3.7, 0.9, size=11, color=WHITE)

b = s.shapes.add_shape(1, Inches(1.5), Inches(5.9), Inches(10.3), Inches(1.2))
b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x03,0x14,0x26); b.line.fill.background()
txt(s, '"Du signal brut à la décision clinique — une brique à la fois."',
    1.5, 6.0, 10.3, 0.9, size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 20 — Questions
# ══════════════════════════════════════════════
s = new_slide()
accent_bar(s, y=3.5, color=ACCENT)
accent_bar(s, y=3.56, color=ACCENT2)

txt(s, "Merci", 1.5, 0.8, 10.3, 1.6, size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Questions ?", 1.5, 2.5, 10.3, 1.0, size=30, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "Dylan Andrade Pereira  ·  Ynov Campus Toulouse  ·  Juin 2026",
    1.5, 4.2, 10.3, 0.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "NeuralES  —  Système d'analyse EEG de la fatigue cognitive",
    1.5, 4.75, 10.3, 0.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)

for i, col in enumerate([ACCENT, ACCENT2, YELLOW, RED]):
    b_deco = s.shapes.add_shape(1, Inches(0.3 + i*0.25), Inches(4.0), Inches(0.08), Inches(2.5))
    b_deco.fill.solid(); b_deco.fill.fore_color.rgb = col; b_deco.line.fill.background()

# ── Sauvegarder ──
out = r"c:\Users\Dylan\Desktop\Side\NeuralES\docs\NeuralES_Presentation.pptx"
prs.save(out)
print(f"OK  Fichier genere : {out}")
print(f"    {len(prs.slides)} slides")
